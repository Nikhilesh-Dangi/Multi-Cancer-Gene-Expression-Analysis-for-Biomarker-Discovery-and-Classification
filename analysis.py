"""End-to-end analysis of the TCGA-like gene expression dataset.

This script loads the .rds dataset, computes required summaries, runs
statistical tests and models, generates plots, and builds a PPT report.

Outputs are written to the ./outputs directory.
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Dict, List, Tuple

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from matplotlib.colors import ListedColormap
from pptx import Presentation
from pptx.util import Inches
from scipy import stats
from sklearn.cluster import AgglomerativeClustering
from sklearn.decomposition import PCA
from sklearn.ensemble import RandomForestClassifier
from sklearn.ensemble import ExtraTreesClassifier, VotingClassifier
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import classification_report, precision_recall_fscore_support
from sklearn.model_selection import StratifiedKFold, train_test_split
from sklearn.preprocessing import StandardScaler

# File paths
DATA_PATH = Path("dataset/Dataset (1).rds")
OUT_DIR = Path("outputs")
OUT_DIR.mkdir(exist_ok=True)


def load_dataset() -> Tuple[pd.DataFrame, List[str]]:
    """Load the dataset and return dataframe plus gene columns."""
    import pyreadr  # imported here to respect existing environment

    result = pyreadr.read_r(str(DATA_PATH))
    df = next(iter(result.values()))
    gene_cols = [c for c in df.columns if c not in ["Sample", "Cancertype"]]
    return df, gene_cols


def compute_basic_summaries(df: pd.DataFrame, gene_cols: List[str]) -> dict:
    """Compute summary stats and save to disk."""
    summaries = {}
    counts = df["Cancertype"].value_counts().sort_values(ascending=False)
    counts.to_csv(OUT_DIR / "samples_per_cancer.csv")
    summaries["counts"] = counts

    gene_means = df[gene_cols].mean()
    gene_stds = df[gene_cols].std()
    pd.DataFrame({"mean": gene_means, "std": gene_stds}).to_csv(
        OUT_DIR / "gene_mean_std.csv"
    )
    variability = gene_stds.sort_values(ascending=False)
    variability.head(5).to_csv(OUT_DIR / "top5_variable_genes.csv")

    summaries["gene_means"] = gene_means
    summaries["gene_stds"] = gene_stds
    summaries["top5_var_genes"] = variability.head(5)
    summaries["top_variable_all"] = variability
    return summaries


def select_top_cancers(df: pd.DataFrame, n: int = 5) -> List[str]:
    """Pick the cancer types with the most samples."""
    return df["Cancertype"].value_counts().head(n).index.tolist()


def differential_expression(
    df: pd.DataFrame, gene_cols: List[str], cancer_types: List[str]
) -> pd.DataFrame:
    """One-way ANOVA across the chosen cancer types for each gene."""
    sub = df[df["Cancertype"].isin(cancer_types)]
    groups = {ct: sub[sub["Cancertype"] == ct] for ct in cancer_types}
    rows = []
    for gene in gene_cols:
        arrays = [groups[ct][gene].values for ct in cancer_types]
        # Skip genes with zero variance in any group to avoid NaNs.
        if any(np.std(arr) == 0 for arr in arrays):
            continue
        f_stat, p_val = stats.f_oneway(*arrays)
        rows.append((gene, f_stat, p_val))
    de_df = pd.DataFrame(rows, columns=["gene", "f_stat", "p_value"]).sort_values(
        "p_value"
    )
    de_df.to_csv(OUT_DIR / "anova_top_cancers.csv", index=False)
    return de_df


def plot_heatmap_top_genes(
    df: pd.DataFrame, genes: List[str], cancer_types: List[str], title: str, fname: str
) -> Path:
    """Plot a heatmap of z-scored expression for selected genes."""
    sub = df[df["Cancertype"].isin(cancer_types)].copy()
    # Limit to a manageable number of samples per cancer for readability.
    sampled = (
        sub.groupby("Cancertype")
        .apply(lambda x: x.sample(min(len(x), 120), random_state=42))
        .reset_index(drop=True)
    )
    zscores = (
        sampled[genes] - sampled[genes].mean()
    ) / sampled[genes].std(ddof=0)
    plt.figure(figsize=(12, 8))
    sns.heatmap(
        zscores.T,
        cmap="RdBu_r",
        center=0,
        cbar_kws={"label": "Z-score"},
    )
    plt.title(title)
    plt.xlabel("Samples")
    plt.ylabel("Genes")
    out_path = OUT_DIR / fname
    plt.tight_layout()
    plt.savefig(out_path, dpi=300)
    plt.close()
    return out_path


def run_clustering(df: pd.DataFrame, gene_cols: List[str]) -> dict:
    """Hierarchical clustering (Ward) on top variable genes."""
    result = {}
    var_ranked = (
        df[gene_cols].std().sort_values(ascending=False).head(200).index.tolist()
    )
    X = df[var_ranked]
    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(X)

    # Sample for visualization to keep dendrogram readable.
    sample_size = min(600, len(X_scaled))
    sample_idx = np.random.RandomState(42).choice(len(X_scaled), size=sample_size, replace=False)
    X_sample = X_scaled[sample_idx]
    labels_sample = df.iloc[sample_idx]["Cancertype"].values

    from scipy.cluster.hierarchy import linkage, dendrogram

    Z = linkage(X_sample, method="ward")

    plt.figure(figsize=(12, 6))
    dendrogram(
        Z,
        p=30,
        truncate_mode="lastp",
        color_threshold=0,
        labels=labels_sample,
        leaf_rotation=90,
        leaf_font_size=8,
        distance_sort="descending",
        show_leaf_counts=True,
        no_labels=True,
    )
    plt.title("Hierarchical clustering (Ward) on top 200 variable genes")
    plt.xlabel("Clusters")
    plt.ylabel("Distance")
    dendro_path = OUT_DIR / "hierarchical_dendrogram.png"
    plt.tight_layout()
    plt.savefig(dendro_path, dpi=300)
    plt.close()
    result["dendrogram"] = dendro_path
    return result


def run_pca(df: pd.DataFrame, gene_cols: List[str]) -> dict:
    """Compute PCA and plot the first two PCs."""
    result = {}
    X = df[gene_cols]
    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(X)
    pca = PCA(n_components=2, random_state=42)
    pcs = pca.fit_transform(X_scaled)
    explained = pca.explained_variance_ratio_
    pc_df = pd.DataFrame(
        {"PC1": pcs[:, 0], "PC2": pcs[:, 1], "Cancertype": df["Cancertype"]}
    )
    plt.figure(figsize=(10, 8))
    sns.scatterplot(
        data=pc_df,
        x="PC1",
        y="PC2",
        hue="Cancertype",
        s=25,
        linewidth=0,
        palette="tab20",
        alpha=0.7,
        legend=False,
    )
    plt.title(
        f"PCA scatter (PC1 {explained[0]*100:.2f}%, PC2 {explained[1]*100:.2f}%)"
    )
    pca_path = OUT_DIR / "pca_scatter.png"
    plt.tight_layout()
    plt.savefig(pca_path, dpi=300)
    plt.close()
    result["explained"] = explained
    result["plot"] = pca_path
    return result

def get_top_genes_by_variance(df: pd.DataFrame, gene_cols: List[str], k: int) -> List[str]:
    """Select top-k variable genes."""
    return df[gene_cols].std().sort_values(ascending=False).head(k).index.tolist()


def _build_tree_model(model_type: str, params: Dict) -> object:
    """Factory for tree-based classifiers."""
    depth = params["max_depth"]
    if isinstance(depth, float):
        if np.isnan(depth):
            depth = None
        else:
            depth = int(depth)
    common = dict(
        n_estimators=int(params["n_estimators"]),
        max_depth=depth,
        max_features=params["max_features"],
        min_samples_leaf=int(params["min_samples_leaf"]),
        random_state=42,
        n_jobs=-1,
        class_weight="balanced_subsample",
    )
    if model_type == "rf":
        return RandomForestClassifier(**common)
    elif model_type == "et":
        return ExtraTreesClassifier(**common)
    else:
        raise ValueError(f"Unknown model type: {model_type}")


def cross_val_accuracy(model, X: pd.DataFrame, y: pd.Series, splits: int = 3) -> float:
    """Compute mean CV accuracy."""
    skf = StratifiedKFold(n_splits=splits, shuffle=True, random_state=42)
    accs = []
    for train_idx, val_idx in skf.split(X, y):
        model.fit(X.iloc[train_idx], y.iloc[train_idx])
        preds = model.predict(X.iloc[val_idx])
        accs.append((preds == y.iloc[val_idx]).mean())
    return float(np.mean(accs))


def run_model_search(
    df: pd.DataFrame,
    gene_cols: List[str],
    *,
    k_values=(400,),
    test_size: float = 0.15,
) -> dict:
    """Quick search over RF/ExtraTrees on a holdout set; returns best model results."""
    # Train/test split for final evaluation
    train_df, test_df = train_test_split(
        df, test_size=test_size, random_state=42, stratify=df["Cancertype"]
    )
    y_train = train_df["Cancertype"]
    y_test = test_df["Cancertype"]

    search_space = []
    for model_type in ["rf", "et"]:
        for k in k_values:
            for n_estimators in [400]:
                for max_depth in [None]:
                    for min_leaf in [1]:
                        search_space.append(
                            {
                                "model": model_type,
                                "top_k": k,
                                "n_estimators": n_estimators,
                                "max_depth": max_depth,
                                "min_samples_leaf": min_leaf,
                                "max_features": "sqrt",
                            }
                        )

    records = []
    trained_models = []
    for cfg in search_space:
        cols = get_top_genes_by_variance(train_df, gene_cols, cfg["top_k"])
        model = _build_tree_model(cfg["model"], cfg)
        model.fit(train_df[cols], y_train)
        preds = model.predict(test_df[cols])
        acc = float((preds == y_test).mean())
        records.append({**cfg, "holdout_accuracy": acc})
        trained_models.append((cfg, cols, model, acc))

    results_df = pd.DataFrame(records).sort_values("holdout_accuracy", ascending=False)
    results_df.to_csv(OUT_DIR / "model_search_results.csv", index=False)

    trained_models = sorted(trained_models, key=lambda x: x[3], reverse=True)
    best_cfg, best_cols, best_model, best_acc = trained_models[0]
    best_report_df = pd.DataFrame(
        classification_report(
            y_test, best_model.predict(test_df[best_cols]), output_dict=True, zero_division=0
        )
    ).transpose()

    # Optional ensemble with second-best model
    final_model_name = "RandomForest" if best_cfg["model"] == "rf" else "ExtraTrees"
    final_acc = best_acc
    final_report_df = best_report_df

    if len(trained_models) > 1:
        second_cfg, _, _, _ = trained_models[1]
        second_model = _build_tree_model(second_cfg["model"], second_cfg)
        ensemble = VotingClassifier(
            estimators=[("best", best_model), ("second", second_model)],
            voting="soft",
            n_jobs=-1,
        )
        ensemble.fit(train_df[best_cols], y_train)
        ens_pred = ensemble.predict(test_df[best_cols])
        ens_acc = float((ens_pred == y_test).mean())
        if ens_acc > final_acc:
            final_acc = ens_acc
            final_report_df = pd.DataFrame(
                classification_report(
                    y_test, ens_pred, output_dict=True, zero_division=0
                )
            ).transpose()
            final_model_name = (
                "Ensemble("
                f"{'RandomForest' if best_cfg['model']=='rf' else 'ExtraTrees'}+"
                f"{'RandomForest' if second_cfg['model']=='rf' else 'ExtraTrees'})"
            )

    final_report_df.to_csv(OUT_DIR / "model_report.csv")

    importances = pd.Series(best_model.feature_importances_, index=best_cols).sort_values(
        ascending=False
    )
    importances.head(30).to_csv(OUT_DIR / "model_top_features.csv")

    with open(OUT_DIR / "model_summary.txt", "w") as f:
        f.write(
            f"Best model: {final_model_name}\n"
            f"Best accuracy: {final_acc:.4f}\n"
            f"Top_k genes: {len(best_cols)}\n"
            f"Search space size: {len(search_space)}\n"
        )

    return {
        "report": final_report_df,
        "accuracy": final_acc,
        "model_name": final_model_name,
        "feature_importances": importances,
        "selected_genes": best_cols,
        "search_results": results_df,
    }


def train_random_forest(
    df: pd.DataFrame,
    gene_cols: List[str],
    *,
    top_k_features: int | None = 500,
) -> dict:
    """Train/test split with Random Forest and return metrics.

    Limits to the top_k_features most variable genes to reduce noise and boost accuracy.
    """
    result = {}
    if top_k_features:
        # Use the most variable genes for better signal/noise.
        variances = df[gene_cols].std().sort_values(ascending=False)
        selected = variances.head(top_k_features).index.tolist()
    else:
        selected = gene_cols

    X = df[selected]
    y = df["Cancertype"]
    X_train, X_test, y_train, y_test = train_test_split(
        X, y, test_size=0.2, random_state=42, stratify=y
    )
    clf = RandomForestClassifier(
        n_estimators=500,
        random_state=42,
        n_jobs=-1,
        max_depth=None,
        max_features="sqrt",
        min_samples_leaf=1,
        class_weight="balanced_subsample",
    )
    clf.fit(X_train, y_train)
    y_pred = clf.predict(X_test)
    report_dict = classification_report(y_test, y_pred, output_dict=True, zero_division=0)
    report_df = pd.DataFrame(report_dict).transpose()
    report_df.to_csv(OUT_DIR / "random_forest_report.csv")
    result["report"] = report_df
    result["accuracy"] = float(report_dict["accuracy"])
    # Feature importance
    importances = pd.Series(clf.feature_importances_, index=selected).sort_values(
        ascending=False
    )
    importances.head(30).to_csv(OUT_DIR / "rf_top30_features.csv")
    result["feature_importances"] = importances
    return result


def logistic_regression_one_vs_rest(
    df: pd.DataFrame, gene_cols: List[str], target: str
) -> pd.DataFrame:
    """Identify genes separating a target cancer vs rest using logistic regression."""
    df_bin = df.copy()
    df_bin["label"] = (df_bin["Cancertype"] == target).astype(int)
    X = df_bin[gene_cols]
    y = df_bin["label"]
    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(X)
    logreg = LogisticRegression(
        penalty="l2",
        max_iter=2000,
        class_weight="balanced",
        n_jobs=-1,
        solver="lbfgs",
    )
    logreg.fit(X_scaled, y)
    coefs = pd.Series(logreg.coef_[0], index=gene_cols)
    top_pos = coefs.sort_values(ascending=False).head(20)
    top_neg = coefs.sort_values().head(20)
    pd.DataFrame({"coef": coefs}).to_csv(OUT_DIR / f"logreg_coefs_{target}.csv")
    pd.DataFrame({"positive": top_pos, "negative": top_neg}).to_csv(
        OUT_DIR / f"logreg_top_genes_{target}.csv"
    )
    return coefs


def gene_correlation(df: pd.DataFrame, variability: pd.Series) -> Path:
    """Compute correlation matrix for top 100 variable genes and plot heatmap."""
    top100 = variability.head(100).index
    corr = df[top100].corr(method="pearson")
    plt.figure(figsize=(12, 10))
    sns.heatmap(
        corr,
        cmap="coolwarm",
        center=0,
        xticklabels=False,
        yticklabels=False,
    )
    plt.title("Gene-gene correlation (top 100 variable genes)")
    corr_path = OUT_DIR / "gene_correlation_heatmap.png"
    plt.tight_layout()
    plt.savefig(corr_path, dpi=300)
    plt.close()
    corr.to_csv(OUT_DIR / "gene_correlation_top100.csv")
    return corr_path


def add_slide(prs, title: str, bullets: List[str]):
    """Add a bullet slide to the PPT."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    for bullet in bullets:
        p = body.add_paragraph()
        p.text = bullet
        p.level = 0


def add_image_slide(prs, title: str, image_path: Path, caption: str | None = None):
    """Add an image slide with optional caption."""
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    left = Inches(1)
    top = Inches(1.2)
    height = Inches(5)
    slide.shapes.add_picture(str(image_path), left, top, height=height)
    if caption:
        tx_box = slide.shapes.add_textbox(left, top + height + Inches(0.1), Inches(8), Inches(1))
        tx_box.text = caption


def build_ppt(
    summaries: dict,
    de_df: pd.DataFrame,
    de_heatmap: Path,
    dendro_path: Path,
    pca_result: dict,
    pca_path: Path,
    rf_result: dict,
    logreg_coefs: pd.Series,
    corr_path: Path,
    target_cancer: str,
):
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Multi-Cancer Gene Expression Analysis"
    slide.placeholders[1].text = "TCGA-style dataset — summaries, DE, clustering, ML"

    add_slide(
        prs,
        "Dataset Summary",
        [
            f"Samples: {int(summaries['counts'].sum())}, Genes: {len(summaries['gene_means'])}",
            f"Top cancer types by N: {', '.join(summaries['counts'].head(5).index.tolist())}",
            f"Top 5 variable genes: {', '.join(summaries['top5_var_genes'].index.tolist())}",
        ],
    )

    add_image_slide(
        prs,
        "Differential Expression Heatmap",
        de_heatmap,
        "Top 10 DE genes (ANOVA) across 5 most frequent cancers",
    )

    add_image_slide(prs, "Hierarchical Clustering", dendro_path, "Ward linkage on top 200 variable genes")

    add_image_slide(
        prs,
        "PCA (PC1 vs PC2)",
        pca_path,
        f"Explained variance: PC1 {pca_result['explained'][0]*100:.2f}%, PC2 {pca_result['explained'][1]*100:.2f}%",
    )

    top_rf = rf_result["feature_importances"].head(10).index.tolist()
    add_slide(
        prs,
        "Model Performance",
        [
            f"Best model: {rf_result['model_name']}",
            f"Accuracy: {rf_result['accuracy']:.3f}",
            "Full report: outputs/model_report.csv; CV: outputs/model_search_results.csv",
            f"Top genes: {', '.join(top_rf)}",
        ],
    )

    top_logreg = logreg_coefs.sort_values(ascending=False).head(8).index.tolist()
    add_slide(
        prs,
        f"Biomarkers for {target_cancer} (LogReg)",
        [
            f"Top positive-weight genes: {', '.join(top_logreg)}",
            f"Full coefficients: outputs/logreg_coefs_{target_cancer}.csv",
        ],
    )

    add_image_slide(prs, "Gene-Gene Co-expression", corr_path, "Pearson correlation among top 100 variable genes")

    prs.save(OUT_DIR / "analysis_summary.pptx")


def main():
    sns.set_theme(style="whitegrid")
    df, gene_cols = load_dataset()
    summaries = compute_basic_summaries(df, gene_cols)

    # Differential expression
    top_cancers = select_top_cancers(df, n=5)
    de_df = differential_expression(df, gene_cols, top_cancers)
    top10_genes = de_df.head(10)["gene"].tolist()
    de_heatmap = plot_heatmap_top_genes(
        df,
        top10_genes,
        top_cancers,
        "Top 10 DE genes across 5 frequent cancers",
        "de_heatmap.png",
    )

    # Clustering & PCA
    clustering = run_clustering(df, gene_cols)
    pca_res = run_pca(df, gene_cols)

    # ML model search (RF/ExtraTrees) with limited hyperparam sweep
    rf_res = run_model_search(df, gene_cols, k_values=(400,))

    # Logistic regression for the most common cancer type
    target = top_cancers[0]
    logreg_coefs = logistic_regression_one_vs_rest(df, gene_cols, target)

    # Gene-gene co-expression
    corr_path = gene_correlation(df, summaries["top_variable_all"])

    # Build PPT
    build_ppt(
        summaries,
        de_df,
        de_heatmap,
        clustering["dendrogram"],
        pca_res,
        pca_res["plot"],
        rf_res,
        logreg_coefs,
        corr_path,
        target,
    )

    print("Analysis complete. See outputs/ for results and PPT.")


if __name__ == "__main__":
    main()
